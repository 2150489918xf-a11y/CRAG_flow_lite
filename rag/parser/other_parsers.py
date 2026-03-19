"""
TXT / HTML / JSON / PPT 解析器
轻量级实现
"""
import logging
import re
import json as json_lib

logger = logging.getLogger(__name__)


# ==================== TXT ====================

def parse_txt(filename, binary=None):
    """解析纯文本文件"""
    try:
        if binary:
            from rag.nlp import find_codec
            codec = find_codec(binary)
            text = binary.decode(codec)
        else:
            with open(filename, "r", encoding="utf-8") as f:
                text = f.read()
    except Exception as e:
        logger.error(f"Failed to read TXT {filename}: {e}")
        return [], []

    sections = []
    # 按空行拆分段落
    paragraphs = re.split(r"\n\s*\n", text)
    for para in paragraphs:
        para = para.strip()
        if para:
            sections.append((para, "text"))

    return sections, []


# ==================== HTML ====================

def parse_html(filename, binary=None):
    """解析 HTML 文件，提取正文和表格"""
    from bs4 import BeautifulSoup

    try:
        if binary:
            from rag.nlp import find_codec
            codec = find_codec(binary)
            html = binary.decode(codec)
        else:
            with open(filename, "r", encoding="utf-8") as f:
                html = f.read()
    except Exception as e:
        logger.error(f"Failed to read HTML {filename}: {e}")
        return [], []

    soup = BeautifulSoup(html, "html.parser")

    # 移除 script 和 style
    for tag in soup.find_all(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    sections = []
    tables = []

    # 提取表格
    for table in soup.find_all("table"):
        tables.append(str(table))
        table.decompose()

    # 提取正文段落
    for elem in soup.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6", "li", "div"]):
        text = elem.get_text(strip=True)
        if text and len(text) > 1:
            tag = "title" if elem.name in ["h1", "h2", "h3"] else "text"
            sections.append((text, tag))

    # 如果段落提取结果很少，fallback 到全文本
    if not sections:
        text = soup.get_text()
        paragraphs = re.split(r"\n\s*\n", text)
        for para in paragraphs:
            para = para.strip()
            if para:
                sections.append((para, "text"))

    return sections, tables


# ==================== JSON ====================

def parse_json(filename, binary=None):
    """解析 JSON 文件，扁平化为文本段落"""
    try:
        if binary:
            from rag.nlp import find_codec
            codec = find_codec(binary)
            text = binary.decode(codec)
        else:
            with open(filename, "r", encoding="utf-8") as f:
                text = f.read()
        data = json_lib.loads(text)
    except Exception as e:
        logger.error(f"Failed to read JSON {filename}: {e}")
        return [], []

    sections = []

    def flatten(obj, prefix=""):
        if isinstance(obj, dict):
            for k, v in obj.items():
                key = f"{prefix}.{k}" if prefix else k
                flatten(v, key)
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                flatten(item, f"{prefix}[{i}]")
        else:
            text = f"{prefix}: {obj}" if prefix else str(obj)
            if text.strip():
                sections.append((text, "text"))

    flatten(data)
    return sections, []


# ==================== PPT ====================

def parse_ppt(filename, binary=None):
    """解析 PPT 文件"""
    from pptx import Presentation
    import io

    try:
        if binary:
            prs = Presentation(io.BytesIO(binary))
        else:
            prs = Presentation(filename)
    except Exception as e:
        logger.error(f"Failed to open PPT {filename}: {e}")
        return [], []

    sections = []
    tables = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    header = rows[0]
                    html = "<table>\n<tr>" + "".join(f"<th>{h}</th>" for h in header) + "</tr>\n"
                    for row in rows[1:]:
                        html += "<tr>" + "".join(f"<td>{c}</td>" for c in row) + "</tr>\n"
                    html += "</table>"
                    tables.append(html)

        if slide_texts:
            sections.append(("\n".join(slide_texts), f"slide_{slide_num}"))

    return sections, tables
